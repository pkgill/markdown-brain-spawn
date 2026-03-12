"""OCR module: extracts text from PDFs, images, and Word documents.

Strategy:
  PDF   → try pdfplumber (embedded text); if sparse, fall back to
          PyMuPDF page rendering + Tesseract per page.
  Image → preprocess with Pillow, then Tesseract.
  DOCX  → python-docx (pure Python; paragraphs + tables).
  PPTX  → python-pptx (pure Python).
  HTML  → BeautifulSoup4 (extracts visible text).
  TXT/SVG → Direct file read (UTF-8).
  DOC/PPT → LibreOffice --headless conversion to text (requires LibreOffice
          or antiword/catppt to be installed on the host system).
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif", ".webp", ".heic", ".heif"}
SUPPORTED_PDF_EXTENSIONS = {".pdf"}
SUPPORTED_WORD_EXTENSIONS = {".doc", ".docx"}
SUPPORTED_POWERPOINT_EXTENSIONS = {".ppt", ".pptx"}
SUPPORTED_TEXT_EXTENSIONS = {".txt", ".svg", ".html", ".htm"}
SUPPORTED_EXTENSIONS = (
    SUPPORTED_IMAGE_EXTENSIONS | 
    SUPPORTED_PDF_EXTENSIONS | 
    SUPPORTED_WORD_EXTENSIONS | 
    SUPPORTED_POWERPOINT_EXTENSIONS | 
    SUPPORTED_TEXT_EXTENSIONS
)


# ---------------------------------------------------------------------------
# Public result type
# ---------------------------------------------------------------------------

@dataclass
class OCRResult:
    text: str
    method: str          # "pdfplumber" | "tesseract_pdf" | "tesseract_image"
    page_count: int | None = None
    confidence: float | None = None   # Tesseract mean confidence 0-100 if available


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def extract_text(path: Path, min_embedded_chars: int = 100,
                 pdf_render_dpi: int = 300, language: str = "eng") -> OCRResult:
    """Extract text from a PDF or image file.

    Args:
        path: Path to the file.
        min_embedded_chars: Minimum chars from pdfplumber before triggering Tesseract fallback.
        pdf_render_dpi: DPI for rendering PDF pages to images (Tesseract fallback).
        language: Tesseract language string, e.g. "eng" or "eng+fra".

    Returns:
        OCRResult with extracted text and method metadata.

    Raises:
        OCRError: If the file cannot be processed.
    """
    suffix = path.suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFormatError(
            f"Unsupported file type: {suffix!r}. "
            f"Supported: {sorted(SUPPORTED_EXTENSIONS)}"
        )

    if not path.exists():
        raise OCRError(f"File not found: {path}")

    if suffix in SUPPORTED_PDF_EXTENSIONS:
        return _extract_pdf(path, min_embedded_chars, pdf_render_dpi, language)
    if suffix in SUPPORTED_WORD_EXTENSIONS:
        return _extract_word(path)
    if suffix in SUPPORTED_POWERPOINT_EXTENSIONS:
        return _extract_powerpoint(path)
    if suffix in SUPPORTED_TEXT_EXTENSIONS:
        return _extract_text_file(path)
    return _extract_image(path, language)


# ---------------------------------------------------------------------------
# PDF extraction
# ---------------------------------------------------------------------------

def _extract_pdf(path: Path, min_embedded_chars: int,
                 pdf_render_dpi: int, language: str) -> OCRResult:
    if _pdf_has_embedded_text(path, min_embedded_chars):
        logger.debug("PDF has embedded text, using pdfplumber: %s", path.name)
        return _extract_pdf_with_pdfplumber(path)
    logger.debug("PDF appears to be scanned, using Tesseract: %s", path.name)
    return _extract_pdf_with_tesseract(path, pdf_render_dpi, language)


def _pdf_has_embedded_text(path: Path, min_chars: int) -> bool:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            if not pdf.pages:
                return False
            # Sample first 3 pages
            sample = pdf.pages[:3]
            total_chars = sum(
                len(p.extract_text() or "") for p in sample
            )
            return total_chars >= min_chars
    except Exception as exc:
        logger.warning("pdfplumber check failed for %s: %s", path.name, exc)
        return False


def _extract_pdf_with_pdfplumber(path: Path) -> OCRResult:
    try:
        import pdfplumber
        pages_text: list[str] = []
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                text = page.extract_text() or ""
                pages_text.append(text)

        full_text = "\n\n".join(t for t in pages_text if t.strip())
        return OCRResult(text=full_text.strip(), method="pdfplumber", page_count=page_count)
    except Exception as exc:
        raise OCRError(f"pdfplumber failed on {path.name}: {exc}") from exc


def _extract_pdf_with_tesseract(path: Path, dpi: int, language: str) -> OCRResult:
    _ensure_tesseract()
    try:
        import fitz  # PyMuPDF
        import pytesseract
        from PIL import Image

        pages_text: list[str] = []
        doc = fitz.open(str(path))
        page_count = len(doc)

        mat = fitz.Matrix(dpi / 72, dpi / 72)   # 72 dpi is the PDF base
        for page in doc:
            pix = page.get_pixmap(matrix=mat)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img = _preprocess_image(img)
            text = pytesseract.image_to_string(img, lang=language)
            pages_text.append(text)

        doc.close()
        full_text = "\n\n".join(t for t in pages_text if t.strip())
        return OCRResult(text=full_text.strip(), method="tesseract_pdf", page_count=page_count)
    except OCRError:
        raise
    except Exception as exc:
        raise OCRError(f"Tesseract PDF extraction failed on {path.name}: {exc}") from exc


# ---------------------------------------------------------------------------
# Word document extraction
# ---------------------------------------------------------------------------

def _extract_word(path: Path) -> OCRResult:
    """Extract text from .docx or legacy .doc files.

    DOCX: uses python-docx (pure Python).
    DOC:  tries LibreOffice --headless conversion first, then antiword.
    """
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return _extract_docx(path)
    return _extract_doc(path)


def _extract_docx(path: Path) -> OCRResult:
    """Extract text from a .docx file using python-docx."""
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise OCRError(
            "python-docx is required to process .docx files.\n"
            "Install it with: pip install python-docx"
        ) from exc

    try:
        doc = docx.Document(str(path))
        parts: list[str] = []

        # Body paragraphs
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        # Tables (each cell on its own line, rows separated by pipe)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)

        full_text = "\n".join(parts)
        return OCRResult(
            text=full_text.strip(),
            method="python-docx",
            page_count=None,
        )
    except OCRError:
        raise
    except Exception as exc:
        raise OCRError(f"python-docx failed on {path.name}: {exc}") from exc


def _extract_doc(path: Path) -> OCRResult:
    """Extract text from a legacy .doc file.

    Tries LibreOffice first (most feature-complete), then antiword.
    Both must be installed as system tools; neither is a pure-Python solution.
    """
    import subprocess
    import tempfile
    import shutil

    # --- Attempt 1: LibreOffice headless conversion to text ---
    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if libreoffice:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    [
                        libreoffice,
                        "--headless",
                        "--convert-to", "txt:Text",
                        "--outdir", tmpdir,
                        str(path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=60,
                )
                if result.returncode == 0:
                    txt_file = Path(tmpdir) / (path.stem + ".txt")
                    if txt_file.exists():
                        text = txt_file.read_text(encoding="utf-8", errors="replace")
                        return OCRResult(
                            text=text.strip(),
                            method="libreoffice",
                            page_count=None,
                        )
        except Exception as exc:
            logger.warning("LibreOffice conversion failed for %s: %s", path.name, exc)

    # --- Attempt 2: antiword ---
    antiword = shutil.which("antiword")
    if antiword:
        try:
            result = subprocess.run(
                [antiword, str(path)],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0:
                return OCRResult(
                    text=result.stdout.strip(),
                    method="antiword",
                    page_count=None,
                )
        except Exception as exc:
            logger.warning("antiword failed for %s: %s", path.name, exc)

    raise OCRError(
        f"Cannot extract text from legacy .doc file: {path.name}\n"
        "Install LibreOffice (recommended) or antiword to process .doc files.\n"
        "  Windows: https://www.libreoffice.org/download/libreoffice/\n"
        "  macOS:   brew install libreoffice  OR  brew install antiword\n"
        "  Ubuntu:  sudo apt-get install libreoffice  OR  sudo apt-get install antiword"
    )


# ---------------------------------------------------------------------------
# Image extraction
# ---------------------------------------------------------------------------

def _extract_image(path: Path, language: str) -> OCRResult:
    _ensure_tesseract()
    try:
        import pytesseract
        from PIL import Image

        try:
            from pillow_heif import register_heif_opener
            register_heif_opener()
        except ImportError:
            pass

        img = Image.open(path)
        img = _preprocess_image(img)

        # image_to_data gives per-word confidence; image_to_string gives clean text
        data = pytesseract.image_to_data(
            img, lang=language, output_type=pytesseract.Output.DICT
        )
        text = pytesseract.image_to_string(img, lang=language)

        # Compute mean confidence (filter out -1 sentinel values)
        confidences = [c for c in data["conf"] if isinstance(c, (int, float)) and c >= 0]
        mean_conf = sum(confidences) / len(confidences) if confidences else None

        return OCRResult(
            text=text.strip(),
            method="tesseract_image",
            page_count=1,
            confidence=mean_conf,
        )
    except OCRError:
        raise
    except Exception as exc:
        raise OCRError(f"Tesseract image extraction failed on {path.name}: {exc}") from exc


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _preprocess_image(img: "Image.Image") -> "Image.Image":
    """Convert to grayscale and boost contrast — Tesseract performs better on preprocessed images."""
    from PIL import ImageEnhance
    img = img.convert("L")   # Grayscale
    enhancer = ImageEnhance.Contrast(img)
    img = enhancer.enhance(2.0)
    return img


def _ensure_tesseract() -> None:
    """Raise a helpful OCRError if the Tesseract binary is not available."""
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
    except Exception as exc:
        raise OCRError(
            "Tesseract OCR binary not found or not working.\n"
            "  macOS:   brew install tesseract\n"
            "  Ubuntu:  sudo apt-get install tesseract-ocr\n"
            "  Windows: https://github.com/UB-Mannheim/tesseract/wiki\n"
            "           Then set ocr.tesseract_cmd in config.yaml to the .exe path."
        ) from exc


def configure_tesseract(cmd: str | None) -> None:
    """Set the Tesseract binary path if provided in config (Windows support)."""
    if cmd:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = cmd


# ---------------------------------------------------------------------------
# PowerPoint extraction
# ---------------------------------------------------------------------------

def _extract_powerpoint(path: Path) -> OCRResult:
    """Extract text from .pptx or legacy .ppt files."""
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return _extract_pptx(path)
    return _extract_ppt(path)


def _extract_pptx(path: Path) -> OCRResult:
    """Extract text from a .pptx file using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise OCRError(
            "python-pptx is required to process .pptx files.\n"
            "Install it with: pip install python-pptx"
        ) from exc

    try:
        prs = Presentation(str(path))
        parts: list[str] = []

        for i, slide in enumerate(prs.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            
            if slide_parts:
                parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_parts))

        full_text = "\n\n".join(parts)
        return OCRResult(
            text=full_text.strip(),
            method="python-pptx",
            page_count=len(prs.slides),
        )
    except Exception as exc:
        raise OCRError(f"python-pptx failed on {path.name}: {exc}") from exc


def _extract_ppt(path: Path) -> OCRResult:
    """Extract text from a legacy .ppt file using LibreOffice."""
    import subprocess
    import tempfile
    import shutil

    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if libreoffice:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    [
                        libreoffice,
                        "--headless",
                        "--convert-to", "txt:Text",
                        "--outdir", tmpdir,
                        str(path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=60,
                )
                if result.returncode == 0:
                    txt_file = Path(tmpdir) / (path.stem + ".txt")
                    if txt_file.exists():
                        text = txt_file.read_text(encoding="utf-8", errors="replace")
                        return OCRResult(
                            text=text.strip(),
                            method="libreoffice_ppt",
                            page_count=None,
                        )
        except Exception as exc:
            logger.warning("LibreOffice PPT conversion failed for %s: %s", path.name, exc)

    raise OCRError(
        f"Cannot extract text from legacy .ppt file: {path.name}\n"
        "Install LibreOffice to process .ppt files."
    )


# ---------------------------------------------------------------------------
# Text/Web/SVG extraction
# ---------------------------------------------------------------------------

def _extract_text_file(path: Path) -> OCRResult:
    """Extract text from .txt, .svg, or .html files."""
    suffix = path.suffix.lower()
    
    # SVG and TXT are direct reads
    if suffix in {".txt", ".svg"}:
        try:
            content = path.read_text(encoding="utf-8", errors="replace")
            return OCRResult(text=content.strip(), method="direct_read")
        except Exception as exc:
            raise OCRError(f"Failed to read text file {path.name}: {exc}") from exc

    # HTML requires BeautifulSoup to extract visible text
    if suffix in {".html", ".htm"}:
        try:
            from bs4 import BeautifulSoup
            content = path.read_text(encoding="utf-8", errors="replace")
            soup = BeautifulSoup(content, "html.parser")
            
            # Remove script and style elements
            for script_or_style in soup(["script", "style"]):
                script_or_style.decompose()

            # Get text, collapse whitespace
            text = soup.get_text(separator="\n")
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = "\n".join(chunk for chunk in chunks if chunk)

            return OCRResult(text=text.strip(), method="beautifulsoup4")
        except ImportError:
            # Fallback to direct read if BS4 is missing (less ideal but better than failing)
            content = path.read_text(encoding="utf-8", errors="replace")
            return OCRResult(text=content.strip(), method="direct_read_html_fallback")
        except Exception as exc:
            raise OCRError(f"Failed to parse HTML {path.name}: {exc}") from exc

    raise OCRError(f"Unexpected text extension: {suffix}")


# ---------------------------------------------------------------------------
# Exception
# ---------------------------------------------------------------------------

class OCRError(Exception):
    """Raised when OCR extraction fails for any reason."""


class UnsupportedFormatError(OCRError):
    """Raised specifically when a file format is not supported by the pipeline."""
